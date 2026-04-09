#!/usr/bin/env python3
"""
v2 Build: DOM-native PPTX reconstruction.

Input: /tmp/pptx-work/v2/{day1,day2}.json, /tmp/pptx-work/v2/assets/*.png
Output: ~/Desktop/서울글로벌센터/{day1_1to4교시,day2_5to8교시}.pptx

Maps each HTML element to a native PPTX shape, text box, or rasterized picture.
Gradients are applied via raw XML patch since python-pptx doesn't expose them.
"""
import json
import os
import re
from copy import deepcopy
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn, nsmap

SLIDE_W_PX = 1280
SLIDE_H_PX = 720
SLIDE_W_EMU = 12192000  # 13.333"
SLIDE_H_EMU = 6858000   # 7.5"
PX_TO_EMU = SLIDE_W_EMU / SLIDE_W_PX  # 9525

DEFAULT_FONT = "Apple SD Gothic Neo"
DEFAULT_COLOR = RGBColor(0x4A, 0x3A, 0x2E)

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def set_run_font(run, font_name):
    """Set latin + east-asian + complex-script font so Korean renders correctly."""
    rPr = run._r.get_or_add_rPr()
    # Remove existing latin/ea/cs
    for t in ('latin', 'ea', 'cs'):
        for el in rPr.findall(qn(f'a:{t}')):
            rPr.remove(el)
    for t in ('latin', 'ea', 'cs'):
        el = etree.SubElement(rPr, qn(f'a:{t}'))
        el.set('typeface', font_name)


def px(v):
    return int(round(v * PX_TO_EMU))


def pt(v):
    """Convert CSS px to PPTX pt with size-based scaling.
    Small fonts get a boost so they remain readable in PPT — Keynote/PowerPoint
    text feels relatively smaller than CSS at the same px because the slide
    canvas is much larger than a browser viewport in physical space.
    """
    px = float(v)
    if px < 16:
        scale = 1.30
    elif px < 20:
        scale = 1.20
    elif px < 26:
        scale = 1.10
    else:
        scale = 1.0
    return Pt(max(6.0, round(px * 0.75 * scale, 2)))


_rgb_re = re.compile(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)(?:\s*,\s*([\d.]+))?\s*\)", re.I)
_hex_re = re.compile(r"#([0-9a-f]{3,8})", re.I)


def parse_rgb(css, default=DEFAULT_COLOR):
    if not css:
        return default, 1.0
    css = css.strip()
    if css == "transparent" or "rgba(0, 0, 0, 0)" in css:
        return default, 0.0
    m = _rgb_re.match(css)
    if m:
        r, g, b = int(m.group(1)), int(m.group(2)), int(m.group(3))
        a = float(m.group(4)) if m.group(4) else 1.0
        return RGBColor(r, g, b), a
    m = _hex_re.match(css)
    if m:
        h = m.group(1)
        if len(h) == 3:
            r, g, b = [int(c * 2, 16) for c in h]
        elif len(h) in (6, 8):
            r = int(h[0:2], 16)
            g = int(h[2:4], 16)
            b = int(h[4:6], 16)
        else:
            return default, 1.0
        return RGBColor(r, g, b), 1.0
    return default, 1.0


def css_align(css):
    if css in ("center", "-webkit-center"):
        return PP_ALIGN.CENTER
    if css == "right":
        return PP_ALIGN.RIGHT
    if css in ("justify", "justify-all"):
        return PP_ALIGN.JUSTIFY
    return PP_ALIGN.LEFT


def set_solid_fill(shape, rgb, alpha=1.0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    if alpha < 1.0:
        # patch alpha
        sp = shape.fill.fore_color._xFill
        srgb = sp.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha_el = etree.SubElement(srgb, qn('a:alpha'))
            alpha_el.set('val', str(int(alpha * 100000)))


def set_gradient_fill(shape, gradient_data):
    """Apply a linear gradient by replacing the <a:solidFill> / <a:noFill> in spPr."""
    sp = shape.fill._xPr  # <p:spPr>
    # Remove any existing fill children
    for tag in ['a:solidFill', 'a:noFill', 'a:gradFill', 'a:blipFill', 'a:pattFill']:
        for el in sp.findall(qn(tag)):
            sp.remove(el)

    grad_fill = etree.SubElement(sp, qn('a:gradFill'))
    grad_fill.set('flip', 'none')
    grad_fill.set('rotWithShape', '1')

    gs_lst = etree.SubElement(grad_fill, qn('a:gsLst'))
    stops = gradient_data.get('stops', [])
    if not stops:
        stops = [
            {'color': 'rgb(200,200,200)', 'pos': 0},
            {'color': 'rgb(100,100,100)', 'pos': 1},
        ]
    for stop in stops:
        gs = etree.SubElement(gs_lst, qn('a:gs'))
        gs.set('pos', str(int(float(stop.get('pos', 0)) * 100000)))
        rgb, alpha = parse_rgb(stop['color'])
        srgb = etree.SubElement(gs, qn('a:srgbClr'))
        srgb.set('val', '{:02X}{:02X}{:02X}'.format(rgb[0], rgb[1], rgb[2]))

    # CSS linear-gradient angle: 0deg = up, 90deg = right, 180deg = down, etc.
    # PPTX gradFill lin ang: 0 = left->right, 90*60000 = up->down
    # CSS angle 180 = top→bottom = PPTX 5400000
    css_angle = gradient_data.get('angle', 180)
    # Convert CSS (clockwise from top=0) to PPTX (clockwise from 3 o'clock=0)
    # CSS 0 = bottom to top = PPTX 270 = 270*60000 = 16200000
    pptx_angle = (css_angle - 90) % 360
    lin = etree.SubElement(grad_fill, qn('a:lin'))
    lin.set('ang', str(int(pptx_angle * 60000)))
    lin.set('scaled', '0')


def set_border(shape, border):
    if not border:
        shape.line.fill.background()
        return
    rgb, alpha = parse_rgb(border['color'])
    t = border.get('top', 0)
    r = border.get('right', 0)
    b = border.get('bottom', 0)
    l = border.get('left', 0)
    sides = [t, r, b, l]
    nonzero = [s for s in sides if s > 0]
    if not nonzero:
        shape.line.fill.background()
        return
    # Symmetric border
    if len(set(sides)) == 1:
        shape.line.color.rgb = rgb
        shape.line.width = px(sides[0])
    else:
        # Asymmetric — skip native border; caller handles with accent rects
        shape.line.fill.background()


def add_border_accents(slide, el):
    """For asymmetric borders, add thin rectangles as accent strips."""
    border = el.get('border')
    if not border:
        return
    bbox = el['bbox']
    t = border.get('top', 0)
    r = border.get('right', 0)
    b = border.get('bottom', 0)
    l = border.get('left', 0)
    if len(set([t, r, b, l])) == 1 or max(t, r, b, l) == 0:
        return  # symmetric, handled via line
    rgb, _ = parse_rgb(border['color'])

    def add_rect(x, y, w, h):
        if w < 0.5 or h < 0.5:
            return
        s = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h)
        )
        s.line.fill.background()
        set_solid_fill(s, rgb)
        s.text_frame.text = ""

    if l > 0:
        add_rect(bbox['x'], bbox['y'], l, bbox['h'])
    if r > 0:
        add_rect(bbox['x'] + bbox['w'] - r, bbox['y'], r, bbox['h'])
    if t > 0:
        add_rect(bbox['x'], bbox['y'], bbox['w'], t)
    if b > 0:
        add_rect(bbox['x'], bbox['y'] + bbox['h'] - b, bbox['w'], b)


def add_shape_element(slide, el):
    bbox = el['bbox']
    if bbox['w'] < 1 or bbox['h'] < 1:
        return

    radius = max(el.get('borderRadius') or [0])
    # Round all 4 corners via shape
    if radius > 0:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
    else:
        shape_type = MSO_SHAPE.RECTANGLE

    shp = slide.shapes.add_shape(
        shape_type,
        px(bbox['x']), px(bbox['y']),
        px(bbox['w']), px(bbox['h']),
    )
    # Remove any default text in the shape
    shp.text_frame.text = ""
    shp.text_frame.margin_left = 0
    shp.text_frame.margin_right = 0
    shp.text_frame.margin_top = 0
    shp.text_frame.margin_bottom = 0

    # Adjust radius
    if radius > 0:
        adj_val = min(0.5, radius / min(bbox['w'], bbox['h']))
        try:
            shp.adjustments[0] = adj_val
        except Exception:
            pass

    # Fill: prefer gradient over solid bg color
    if el.get('gradient') and el['gradient'].get('stops'):
        set_gradient_fill(shp, el['gradient'])
    elif el.get('bg'):
        rgb, alpha = parse_rgb(el['bg'])
        if alpha > 0:
            set_solid_fill(shp, rgb, alpha)
        else:
            shp.fill.background()
    else:
        shp.fill.background()

    set_border(shp, el.get('border'))
    add_border_accents(slide, el)

    # opacity
    op = el.get('opacity', 1.0)
    if op < 1.0:
        # apply as alpha on fill
        try:
            sp = shp.fill.fore_color._xFill
            if sp is not None and sp.find(qn('a:alpha')) is None:
                alpha_el = etree.SubElement(sp, qn('a:alpha'))
                alpha_el.set('val', str(int(op * 100000)))
        except Exception:
            pass

    # box shadow (effects)
    shadow = el.get('boxShadow')
    if shadow:
        try:
            spPr = shp.fill._xPr
            eff_lst = etree.SubElement(spPr, qn('a:effectLst'))
            outer = etree.SubElement(eff_lst, qn('a:outerShdw'))
            outer.set('blurRad', str(px(shadow.get('blur', 0))))
            outer.set('dist', str(px(max(abs(shadow.get('offsetX', 0)), abs(shadow.get('offsetY', 0))))))
            outer.set('dir', '2700000')
            outer.set('algn', 'ctr')
            outer.set('rotWithShape', '0')
            rgb, alpha = parse_rgb(shadow.get('color', 'rgba(0,0,0,0.15)'))
            srgb = etree.SubElement(outer, qn('a:srgbClr'))
            srgb.set('val', '{:02X}{:02X}{:02X}'.format(rgb[0], rgb[1], rgb[2]))
            if alpha < 1.0:
                a_el = etree.SubElement(srgb, qn('a:alpha'))
                a_el.set('val', str(int(alpha * 100000)))
        except Exception:
            pass


def add_slide_background(slide, el):
    """Slide-level background: full-slide shape behind everything."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W_EMU, SLIDE_H_EMU
    )
    shp.text_frame.text = ""
    shp.line.fill.background()
    if el.get('bg') and el['bg'].get('gradient') and el['bg']['gradient'].get('stops'):
        set_gradient_fill(shp, el['bg']['gradient'])
    elif el.get('bg') and el['bg'].get('color'):
        rgb, alpha = parse_rgb(el['bg']['color'])
        if alpha > 0:
            set_solid_fill(shp, rgb, alpha)
        else:
            shp.fill.background()
    else:
        # Default paper background
        set_solid_fill(shp, RGBColor(0xFA, 0xF7, 0xF2))

    # cork-top border: if the slide has borderTopWidth, add a top bar shape
    btw = el.get('borderTopWidth', 0)
    if btw and btw > 0:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W_EMU, px(btw)
        )
        bar.line.fill.background()
        rgb, _ = parse_rgb(el.get('borderTopColor', 'rgb(212,160,96)'))
        set_solid_fill(bar, rgb)


def add_text_element(slide, el):
    runs = el.get('runs', [])
    if not runs:
        return
    # Group runs by line (y coordinate with small tolerance)
    runs_sorted = sorted(runs, key=lambda r: (round(r['y'] / 3), r['x']))
    lines = []
    current_y = None
    current_line = []
    for r in runs_sorted:
        ry = round(r['y'] / 3)
        if current_y is None or abs(ry - current_y) <= 2:
            current_line.append(r)
            current_y = ry
        else:
            lines.append(current_line)
            current_line = [r]
            current_y = ry
    if current_line:
        lines.append(current_line)

    # For each line, create a tight textbox that contains its runs
    for line in lines:
        x0 = min(r['x'] for r in line)
        y0 = min(r['y'] for r in line)
        x1 = max(r['x'] + r['w'] for r in line)
        y1 = max(r['y'] + r['h'] for r in line)
        w = max(x1 - x0, 10)
        h = max(y1 - y0, 12)

        # Slight padding
        tb = slide.shapes.add_textbox(
            px(x0 - 1), px(y0 - 1),
            px(w + 4), px(h + 4),
        )
        tf = tb.text_frame
        tf.word_wrap = False
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.TOP
        try:
            tf.auto_size = None
        except Exception:
            pass

        p = tf.paragraphs[0]
        align = css_align(line[0].get('textAlign', 'left'))
        p.alignment = align

        # Runs in this line
        line_sorted = sorted(line, key=lambda r: r['x'])
        for j, r in enumerate(line_sorted):
            run = p.add_run()
            txt = r.get('text', '')
            # Replace nbsp and clean
            run.text = txt.replace('\u00a0', ' ')
            run.font.size = pt(r.get('fontSize', 16))
            run.font.bold = int(r.get('fontWeight', 400)) >= 600
            run.font.italic = bool(r.get('italic', False))
            rgb, _ = parse_rgb(r.get('color', ''))
            run.font.color.rgb = rgb
            set_run_font(run, DEFAULT_FONT)


def add_raster_element(slide, el, assets_by_id, assets_dir):
    # Find matching asset in assets list
    bbox = el['bbox']
    # Assets are tagged like 'svg-0', 'xfm-0', 'pse-0' per slide
    # Caller passes assets_by_id as dict {tag: png_filename}
    # We need to figure out which tag this element matches.
    # Since the elements and assets were generated in the same DFS order,
    # we use a stateful counter in the caller.
    pass  # handled in build_slide


def build_slide(prs, blank_layout, slide_data, assets_dir):
    slide = prs.slides.add_slide(blank_layout)
    elements = slide_data.get('elements', [])
    assets = slide_data.get('assets', [])

    # Build an iterator of assets by kind (svg, transform, pseudo)
    # Assets are in DFS order, so we consume them sequentially per kind.
    asset_iter_by_kind = {'svg': iter([a for a in assets if a['tag'].startswith('svg')]),
                          'transform': iter([a for a in assets if a['tag'].startswith('xfm')]),
                          'pseudo': iter([a for a in assets if a['tag'].startswith('pse')])}
    svg_i = 0
    xfm_i = 0
    pse_i = 0

    # Emit slide background first
    slide_el = elements[0] if elements and elements[0].get('tag') == 'SLIDE' else None
    if slide_el:
        add_slide_background(slide, slide_el)
        remaining = elements[1:]
    else:
        remaining = elements

    for el in remaining:
        tag = el.get('tag')
        if tag == 'SHAPE':
            add_shape_element(slide, el)
        elif tag == 'TEXT':
            add_text_element(slide, el)
        elif tag == 'RASTER':
            subtype = el.get('subtype', 'svg')
            # Match asset by kind + order
            matching_asset = None
            kind_prefix = {'svg': 'svg-', 'img': 'img-', 'transform': 'xfm-', 'pseudo': 'pse-'}.get(subtype, 'svg-')
            for a in assets:
                if a['tag'].startswith(kind_prefix):
                    matching_asset = a
                    break
            if matching_asset:
                # Remove from list to avoid reusing
                assets.remove(matching_asset)
                png_path = os.path.join(assets_dir, matching_asset['png'])
                if os.path.exists(png_path):
                    bbox = el['bbox']
                    try:
                        slide.shapes.add_picture(
                            png_path,
                            px(bbox['x']), px(bbox['y']),
                            width=px(bbox['w']), height=px(bbox['h']),
                        )
                    except Exception as e:
                        print(f'  WARN picture: {e}')
    return slide


def build(json_path, assets_dir, out_path):
    with open(json_path, 'r', encoding='utf-8') as f:
        slides_data = json.load(f)

    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU
    blank = prs.slide_layouts[6]

    for sd in slides_data:
        build_slide(prs, blank, sd, assets_dir)

    prs.save(out_path)
    print(f'[OK] {out_path}  ({len(slides_data)} slides)')


if __name__ == '__main__':
    base = '/tmp/pptx-work/v2'
    assets = os.path.join(base, 'assets')
    out_dir = os.path.expanduser('~/Desktop/서울글로벌센터')
    os.makedirs(out_dir, exist_ok=True)

    build(
        os.path.join(base, 'day1.json'),
        assets,
        os.path.join(out_dir, 'day1_1to4교시.pptx'),
    )
    build(
        os.path.join(base, 'day2.json'),
        assets,
        os.path.join(out_dir, 'day2_5to8교시.pptx'),
    )
