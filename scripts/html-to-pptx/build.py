#!/usr/bin/env python3
"""
Build editable PPTX from captured HTML slide data.

Input:
  /tmp/pptx-work/day1.json  (per-slide texts + PNG filename)
  /tmp/pptx-work/day2.json
  /tmp/pptx-work/slides/*.png  (text-hidden slide backgrounds)

Output:
  ~/Desktop/서울글로벌센터/day1_1to4교시.pptx
  ~/Desktop/서울글로벌센터/day2_5to8교시.pptx
"""
import json
import os
import re
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SLIDE_W_PX = 1280
SLIDE_H_PX = 720
# 16:9: 13.333" x 7.5" = 12192000 x 6858000 EMU
SLIDE_W_EMU = 12192000
SLIDE_H_EMU = 6858000
PX_TO_EMU = SLIDE_W_EMU / SLIDE_W_PX  # 9525

DEFAULT_FONT = "Noto Sans KR"
DEFAULT_COLOR = RGBColor(0x4A, 0x3A, 0x2E)


def px(v):
    return int(round(v * PX_TO_EMU))


_rgb_re = re.compile(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)")


def parse_color(css: str) -> RGBColor:
    if not css:
        return DEFAULT_COLOR
    m = _rgb_re.match(css.strip())
    if m:
        return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    return DEFAULT_COLOR


def css_align(css: str):
    if css in ("center", "-webkit-center"):
        return PP_ALIGN.CENTER
    if css == "right":
        return PP_ALIGN.RIGHT
    if css in ("justify", "justify-all"):
        return PP_ALIGN.JUSTIFY
    return PP_ALIGN.LEFT


def build(json_path: str, png_dir: str, out_path: str):
    with open(json_path, "r", encoding="utf-8") as f:
        slides = json.load(f)

    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU
    blank = prs.slide_layouts[6]  # Blank

    for sd in slides:
        slide = prs.slides.add_slide(blank)
        # Background image: fill whole slide
        png_path = os.path.join(png_dir, sd["png"])
        if os.path.exists(png_path):
            slide.shapes.add_picture(
                png_path, 0, 0, width=SLIDE_W_EMU, height=SLIDE_H_EMU
            )
        else:
            print(f"  WARN: missing {png_path}")

        for t in sd["texts"]:
            x = max(0, t["x"])
            y = max(0, t["y"])
            w = max(t["w"], 20)
            h = max(t["h"], 14)
            # Clamp to slide
            if x >= SLIDE_W_PX or y >= SLIDE_H_PX:
                continue
            w = min(w, SLIDE_W_PX - x)
            h = min(h, SLIDE_H_PX - y)

            tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
            # Vertical center-ish
            tf.vertical_anchor = MSO_ANCHOR.TOP

            text = t.get("text", "")
            lines = text.split("\n")
            font_pt = max(6.0, round(float(t.get("fontSize", 16)) * 0.75, 1))
            bold = int(t.get("fontWeight", 400)) >= 600
            color = parse_color(t.get("color", ""))
            align = css_align(t.get("textAlign", "left"))

            for i, line in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.alignment = align
                run = p.add_run()
                run.text = line
                run.font.size = Pt(font_pt)
                run.font.name = DEFAULT_FONT
                run.font.bold = bold
                run.font.color.rgb = color

    prs.save(out_path)
    print(f"[OK] {out_path}  ({len(slides)} slides)")


if __name__ == "__main__":
    base = "/tmp/pptx-work"
    png_dir = os.path.join(base, "slides")
    out_dir = os.path.expanduser("~/Desktop/서울글로벌센터")
    os.makedirs(out_dir, exist_ok=True)

    build(
        os.path.join(base, "day1.json"),
        png_dir,
        os.path.join(out_dir, "day1_1to4교시.pptx"),
    )
    build(
        os.path.join(base, "day2.json"),
        png_dir,
        os.path.join(out_dir, "day2_5to8교시.pptx"),
    )
